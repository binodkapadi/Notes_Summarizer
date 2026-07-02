import os
import time
import tempfile
from io import BytesIO
from pathlib import Path
from google import genai
import streamlit as st
import streamlit.components.v1 as components
from bs4 import BeautifulSoup
from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from markdown import markdown
from pptx import Presentation

# Set page config as the very first Streamlit command
st.set_page_config(page_title="Notes Summarizer", page_icon="📝")

# Load environment variables
load_dotenv()

# Configure Gemini API key
client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))

# Mapping between human‑friendly names and Gemini model IDs
MODEL_OPTIONS = {
    "Gemini Flash-Lite Latest": "gemini-flash-lite-latest",
    "Gemini Flash Latest": "gemini-flash-latest",
    "Gemini 2.5 Flash-Lite Preview": "gemini-2.5-flash-lite-preview-09-2025",
    "Gemini 2.5 Flash": "gemini-2.5-flash",
    "Gemini 2.5 Flash-Lite": "gemini-2.5-flash-lite",
}

FONT_CANDIDATES = [
    (
        Path("C:/Windows/Fonts/arial.ttf"),
        Path("C:/Windows/Fonts/arialbd.ttf"),
    ),
    (
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
    ),
    (
        Path("/System/Library/Fonts/Supplemental/Arial.ttf"),
        Path("/System/Library/Fonts/Supplemental/Arial Bold.ttf"),
    ),
    (
        Path("assets/fonts/DejaVuSans.ttf"),
        Path("assets/fonts/DejaVuSans-Bold.ttf"),
    ),
]


def locate_font_paths() -> tuple[Path, Path]:
    """Return paths to regular and bold font files that support Unicode."""
    for regular, bold in FONT_CANDIDATES:
        if regular.exists():
            bold_path = bold if bold and bold.exists() else regular
            return regular, bold_path
    raise FileNotFoundError(
        "Could not locate a Unicode font. Please add a TTF font (e.g. DejaVuSans) "
        "under assets/fonts/ or update FONT_CANDIDATES with the correct paths."
    )


def summarize_notes(all_text: str, gemini_files, model_name: str, has_pdf: bool) -> str:
    """
    Call Gemini to create thorough, human-friendly study notes from text + uploaded files.
    """
    if not all_text and not gemini_files:
        raise ValueError("No content provided to summarize.")

    # The model ID is used directly in the generate_content call

    system_prompt = """
You are a senior AI Engineer, Prompt Engineer, and Full-Stack Developer.

## Objective

Develop an intelligent Notes Generator that creates concise, complete, well-structured study notes from any input source.

The input may be:
* PDF documents
* DOCX documents
* PPT/PPTX presentations
* TXT files
* Markdown files
* Images (after OCR)
* Pasted text
* User prompts
* Multiple uploaded files together

The generated notes must cover **100% of the available content** without becoming lengthy or verbose.

---

# Core Requirements

## 1. Complete Topic Coverage
The generated notes must include every important topic present in the source.
Never skip:
* Main topics
* Subtopics
* Definitions
* Important concepts
* Key points
* Algorithms
* Formulae
* Theorems
* Laws
* Principles
* Examples (only if essential)
* Diagrams (describe briefly if necessary)
* Tables (summarize)
* Comparisons
* Advantages
* Disadvantages
* Applications
* Features
* Characteristics
* Classifications
* Processes
* Workflows
* Important facts
* Conclusions

Every heading and meaningful section from the source should be represented in the notes.

---

# 2. Size and Length Constraints (1/4th Rule)
Do NOT copy paragraphs.
Instead:
* Compress information.
* Remove unnecessary wording.
* Eliminate repetition.
* Keep only meaningful information.
* Use concise sentences.

Target:
* **The generated summarized notes MUST be approximately 1/4th (25%) of the length of the original content.**
* This 1/4th size rule applies to all documents (PDFs, PPTs, DOCs) and any long pasted text or user prompts.
* (Note: The 1/4th length rule does not apply to image-only uploads like JPG/PNG).
* Even with this 25% length constraint, you MUST cover ALL topics and explain each topic clearly. Do not skip topics to save space; instead, compress the explanations intelligently so that all concepts are covered in minimum words.

---

# 3. Preserve Hierarchy
Maintain the logical structure.
Example:
Topic
Definition
Features
Types
Advantages
Disadvantages
Applications
Examples
Conclusion

Subtopics must appear under their parent topic.
Never flatten everything into one list.

---

# 4. Intelligent Compression
Convert long explanations into concise notes.
Example:
Long paragraph
↓
2–5 concise lines
without losing important information.

---

# 5. Important Keywords
Highlight:
* Important terms
* Technical vocabulary
* Definitions
* Formula names
* Algorithms
* Standards
* Protocols

These should stand out clearly.

---

# 6. Include Every Subtopic
Even if a subtopic contains only a few lines, include it.
Never omit sections simply because they are short.

---

# 7. No Hallucinations
Do not invent information.
Generate notes strictly from:
* Uploaded documents
* OCR text
* User-pasted text
* User prompt

If information does not exist, do not create it.

---

# 8. Handle Large Documents
For long PDFs:
Read every page.
Do not summarize only the beginning.
Ensure coverage of:
* First page
* Middle pages
* Last pages
* Appendices (if relevant)

Every page contributes to the final notes.

---

# 9. Merge Duplicate Information
If the same concept appears multiple times:
Merge it into one concise explanation.
Avoid duplication.

---

# 10. Ignore Unnecessary Content
Skip:
* Headers
* Footers
* Page numbers
* Watermarks
* Blank pages
* Decorative elements
* Repeated titles
* Navigation text

unless they contain meaningful academic content.

---

# 11. Maintain Technical Accuracy
Never change:
* Equations
* Formulae
* Units
* Symbols
* Mathematical expressions
* Programming syntax
* Code snippets (summarize only if appropriate)

---

# 12. Tables
If tables contain useful information:
Convert them into concise notes.
Do not reproduce large tables unless necessary.

---

# 13. Figures
If diagrams explain an important concept:
Write a short textual explanation.
Example:
OSI Model
* 7 layers
* Application → Physical
* Data flows downward during transmission

---

# 14. Comparisons
Convert comparisons into compact tables when useful.
Example:
SQL vs NoSQL
Feature | SQL | NoSQL
instead of long paragraphs.

---

# 15. Definitions
Definitions should be:
* One line whenever possible.
* Maximum two lines if necessary.

---

# 16. Lists
Where appropriate use:
* Numbered lists
* Bullet lists
* Small tables

Avoid huge paragraphs.

---

# 17. Study-Friendly Notes
The notes should feel like exam revision notes.
A student should be able to revise quickly.

---

# 18. Preserve Important Order
If the document follows a logical order, maintain that order.
Do not randomly rearrange topics.

---

# 19. Output Format
Use Markdown.
Example:
# Chapter Title
## Topic
Short explanation
### Subtopic
* Point 1
* Point 2
* Point 3

## Next Topic
...

---

# 20. Conciseness Rule
Every sentence must add value.
Remove:
* filler words
* repeated ideas
* unnecessary examples
* storytelling

---

# 21. Coverage Validation
Before producing the final notes, internally verify:
✓ Every chapter included
✓ Every heading covered
✓ Every subheading covered
✓ Every major concept included
✓ No significant topic omitted
✓ No duplicated content
✓ Notes remain concise

Only then generate the final notes.

---

# 22. Quality Goal
The generated notes should:
* Cover all topics from the source.
* Be significantly shorter than the original.
* Be easy to revise before exams.
* Preserve all essential concepts.
* Be accurate and well organized.
* Be readable in a single pass.
* Avoid unnecessary detail while ensuring no important topic or subtopic is missed.

The final output should be concise, comprehensive, structured, exam-oriented, and faithful to the source material.
    """

    contents = [system_prompt]

    if all_text.strip():
        contents.append(all_text.strip())

    # Attach uploaded PDF / image files for multimodal summarization
    contents.extend(gemini_files)

    max_retries = 3
    for attempt in range(max_retries):
        try:
            print(f"[Backend] Calling Gemini API with model {model_name} (Attempt {attempt + 1}/{max_retries})...")
            response = client.models.generate_content(
                model=model_name,
                contents=contents
            )
            print("[Backend] Gemini API generation completed successfully.")        
            return (response.text or "").strip()
        except Exception as e:
            if attempt == max_retries - 1:
                print(f"[Backend] ERROR: Gemini API call failed after {max_retries} attempts. Exception: {e}")
                raise e
            print(f"[Backend] WARNING: Error calling Gemini: {e}. Retrying in {2 ** attempt} seconds...")
            time.sleep(2 ** attempt)


def extract_content_from_docx(file_bytes: bytes) -> tuple[str, list[bytes]]:
    """Extract plain text and images from a Word .docx file."""
    doc = Document(BytesIO(file_bytes))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            images.append(rel.target_part.blob)
            
    return "\n".join(parts), images


def extract_content_from_pptx(file_bytes: bytes) -> tuple[str, list[bytes]]:
    """Extract plain text and images from a PowerPoint .pptx file."""
    prs = Presentation(BytesIO(file_bytes))
    texts = []
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    texts.append(t)
            if hasattr(shape, "image") and shape.image:
                images.append(shape.image.blob)
    return "\n".join(texts), images


def render_footer():
    """Render a footer at the bottom of the page (always visible)."""
    st.markdown(
        """
        <div class="app-footer">
            <p class="footer-text">
                Copyright © by
                <span class="footer-name">Binod Kapadi</span>
            </p>
            <p class="footer-text">All Rights Reserved — 2024</p>
        </div>
        """,
        unsafe_allow_html=True,
    )


def generate_summary_pdf(summary_markdown: str) -> bytes:
    """Convert the summary markdown to a mobile-friendly PDF with UI-like spacing."""
    font_regular, font_bold = locate_font_paths()

    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_text_color(0, 0, 0)
    pdf.add_font("Body", "", str(font_regular), uni=True)
    pdf.add_font("BodyBold", "", str(font_bold), uni=True)

    html = markdown(summary_markdown, extensions=["tables"])
    soup = BeautifulSoup(html, "html.parser")

    def add_spacing(lines: int = 4):
        pdf.ln(lines)

    def line_count(text: str, width: float) -> int:
        """Estimate line count for table cells."""
        words = (text or "").split()
        if not words:
            return 1
        lines = 1
        current = ""
        for word in words:
            tentative = (current + " " if current else "") + word
            if pdf.get_string_width(tentative) > width:
                lines += 1
                current = word
            else:
                current = tentative
        return lines

    first_block = True

    for block in soup.children:
        if not getattr(block, "name", None):
            continue

        tag = block.name.lower()
        text = block.get_text(" ", strip=True)

        if tag in {"h1", "h2", "h3", "h4"}:
            if not first_block:
                add_spacing(6)
            font_sizes = {"h1": 20, "h2": 17, "h3": 15, "h4": 13}
            heading_colors = {
                "h1": (34, 99, 188),
                "h2": (52, 152, 219),
                "h3": (44, 62, 80),
                "h4": (44, 62, 80),
            }
            pdf.set_text_color(*heading_colors.get(tag, (0, 0, 0)))
            pdf.set_font("BodyBold", size=font_sizes.get(tag, 14))
            pdf.multi_cell(0, 8, text)
            pdf.set_text_color(0, 0, 0)
            add_spacing(4)
        elif tag == "p" and text:
            pdf.set_font("Body", size=11)
            pdf.multi_cell(0, 6, text)
            add_spacing(4)
        elif tag in {"ul", "ol"}:
            items = [
                li.get_text(" ", strip=True)
                for li in block.find_all("li", recursive=False)
            ]
            pdf.set_font("Body", size=11)
            add_spacing(2)
            for idx, item in enumerate(items, start=1):
                bullet = "-" if tag == "ul" else f"{idx}."
                pdf.multi_cell(0, 6, f"{bullet} {item}")
            add_spacing(4)
        elif tag == "table":
            add_spacing(4)
            rows = []
            max_cols = 0
            for tr in block.find_all("tr"):
                cells = tr.find_all(["th", "td"])
                row_texts = [c.get_text(" ", strip=True) for c in cells]
                rows.append(row_texts)
                max_cols = max(max_cols, len(row_texts))

            if not rows or max_cols == 0:
                continue

            content_width = pdf.w - pdf.l_margin - pdf.r_margin
            col_width = content_width / max_cols
            line_height = 7

            for idx, row in enumerate(rows):
                normalized = row + [""] * (max_cols - len(row))
                is_header = idx == 0
                pdf.set_font("BodyBold" if is_header else "Body", size=11)
                pdf.set_fill_color(28, 42, 62) if is_header else pdf.set_fill_color(250, 250, 250)
                pdf.set_text_color(255, 255, 255) if is_header else pdf.set_text_color(0, 0, 0)

                start_x = pdf.get_x()
                start_y = pdf.get_y()
                max_lines = 1

                for col_idx, cell_text in enumerate(normalized):
                    lines = max(line_count(cell_text, col_width - 2), 1)
                    max_lines = max(max_lines, lines)
                    pdf.multi_cell(
                        col_width,
                        line_height,
                        cell_text,
                        border=1,
                        align="L",
                        fill=True,
                    )
                    pdf.set_xy(start_x + col_width * (col_idx + 1), start_y)

                pdf.set_xy(start_x, start_y + max_lines * line_height)
                pdf.set_text_color(0, 0, 0)
            add_spacing(6)
        elif tag == "hr":
            pdf.set_draw_color(200, 200, 200)
            y = pdf.get_y()
            pdf.line(10, y, 200, y)
            add_spacing(4)

        first_block = False

    # FPDF returns latin-1 encoded str; ensure bytes output
    pdf_str = pdf.output(dest="S")
    return pdf_str.encode("latin-1", "ignore")


def main():
    st.title("📝 Notes Summarizer")

    # Optional dark mode CSS (force UTF-8 to avoid Windows cp1252 decode errors)
    if os.path.exists("style.css"):
        with open("style.css", "r", encoding="utf-8", errors="ignore") as f:
            st.markdown(f"<style>{f.read()}</style>", unsafe_allow_html=True)

    # Check existing uploaded files in session to decide whether to disable text input
    uploaded_files_state = st.session_state.get("file_uploader")

    # Text input (shown first; disabled when any files are already uploaded)
    text_content = st.text_area(
        "📘 Paste your notes here :",
        height=200,
        placeholder="Paste long notes, lecture text, or copied content here...",
        key="notes_text",
        disabled=bool(uploaded_files_state),
    )

    # After reading text, decide whether to disable file uploader
    text_has_content = bool(text_content.strip())

    # File upload: PDFs, Word docs, PowerPoints, text files, and images
    uploaded_file = st.file_uploader(
        "📂 Or upload notes as PDF, Word, PPT, text, or images (limit 50MB per file):",
        type=["pdf", "doc", "docx", "ppt", "pptx", "txt", "png", "jpg", "jpeg"],
        accept_multiple_files=False,
        help="Limit 50MB per file • Supports PDF, DOC/DOCX, PPT/PPTX, TXT, PNG, JPG.",
        key="file_uploader",
        disabled=text_has_content,
    )

    # Initialize session state for summary
    if "summary" not in st.session_state:
        st.session_state.summary = ""

    # Placeholder for full-page loading overlay
    loader_placeholder = st.empty()

    # Model selection moved just above the Summarize button (no settings icon)
    model_label = st.selectbox("🤖 Select Gemini model:", list(MODEL_OPTIONS.keys()))
    selected_model = MODEL_OPTIONS[model_label]

    if st.button("✨ Summarize"):
        if not text_content.strip() and not uploaded_file:
            st.warning("⚠️ Please provide some notes or upload at least one file.")
        else:
            # Clear previous summary so old notes disappear while new summary is being generated
            st.session_state.summary = ""
            # Show custom blurred overlay loader
            loader_html = """
            <div class="overlay-loader">
                <div class="overlay-loader-circle"></div>
                <p class="overlay-loader-text">Summarizing your notes... Please wait ⏳</p>
            </div>
            """
            loader_placeholder.markdown(loader_html, unsafe_allow_html=True)

            with st.spinner("Summarizing your notes... Please wait ⏳"):
                text_parts = []
                gemini_files = []
                has_pdf = False

                # Include pasted text
                if text_content.strip():
                    text_parts.append(text_content.strip())

                # Handle uploaded files
                max_bytes = 50 * 1024 * 1024  # 50 MB per file limit
                uploads_to_process = [uploaded_file] if uploaded_file else []
                for uploaded in uploads_to_process:
                    suffix = Path(uploaded.name).suffix.lower()

                    # Enforce 50 MB limit for each uploaded file
                    if uploaded.size > max_bytes:
                        loader_placeholder.empty()
                        st.error(f"Error: Selected file '{uploaded.name}' must be less than 50 MB.")
                        st.stop()

                    # Plain text files: read locally and append to text
                    if suffix == ".txt":
                        try:
                            file_text = uploaded.read().decode("utf-8", errors="ignore")
                            text_parts.append(file_text)
                        finally:
                            uploaded.seek(0)

                    # Word documents: extract text and images locally
                    elif suffix in [".doc", ".docx"]:
                        try:
                            file_bytes = uploaded.read()
                            doc_text, doc_images = extract_content_from_docx(file_bytes)
                            if doc_text:
                                text_parts.append(doc_text)
                            for img_bytes in doc_images:
                                fd, temp_path = tempfile.mkstemp(suffix=".png")
                                os.close(fd)
                                try:
                                    print(f"[Backend] Uploading extracted image from '{uploaded.name}' to Gemini...")
                                    with open(temp_path, "wb") as tmp_file:
                                        tmp_file.write(img_bytes)
                                    gemini_file = client.files.upload(file=temp_path)
                                    gemini_files.append(gemini_file)
                                    print(f"[Backend] Successfully uploaded image as {gemini_file.name}")
                                finally:
                                    try:
                                        os.remove(temp_path)
                                    except OSError:
                                        pass
                        except Exception as doc_err:
                            st.warning(f"Could not read Word document '{uploaded.name}': {doc_err}")
                        finally:
                            uploaded.seek(0)

                    # PowerPoint presentations: extract text and images locally
                    elif suffix in [".ppt", ".pptx"]:
                        try:
                            file_bytes = uploaded.read()
                            ppt_text, ppt_images = extract_content_from_pptx(file_bytes)
                            if ppt_text:
                                text_parts.append(ppt_text)
                            for img_bytes in ppt_images:
                                fd, temp_path = tempfile.mkstemp(suffix=".png")
                                os.close(fd)
                                try:
                                    print(f"[Backend] Uploading extracted image from '{uploaded.name}' to Gemini...")
                                    with open(temp_path, "wb") as tmp_file:
                                        tmp_file.write(img_bytes)
                                    gemini_file = client.files.upload(file=temp_path)
                                    gemini_files.append(gemini_file)
                                    print(f"[Backend] Successfully uploaded image as {gemini_file.name}")
                                finally:
                                    try:
                                        os.remove(temp_path)
                                    except OSError:
                                        pass
                        except Exception as ppt_err:
                            st.warning(f"Could not read PowerPoint file '{uploaded.name}': {ppt_err}")
                        finally:
                            uploaded.seek(0)

                    else:
                        # PDFs and images: upload to Gemini for multimodal summarization.
                        # On Windows, NamedTemporaryFile can cause permission issues because the
                        # file handle stays open. Instead, create a temp path, close the handle,
                        # then upload and finally delete the temp file.
                        if suffix == ".pdf":
                            has_pdf = True

                        fd, temp_path = tempfile.mkstemp(suffix=suffix)
                        os.close(fd)
                        try:
                            print(f"[Backend] Uploading '{uploaded.name}' to Gemini...")
                            with open(temp_path, "wb") as tmp_file:
                                tmp_file.write(uploaded.getbuffer())

                            gemini_file = client.files.upload(file=temp_path)
                            
                            # Wait for processing (essential for large PDFs)
                            state = getattr(gemini_file.state, "name", gemini_file.state) if hasattr(gemini_file, "state") else None
                            if state == "PROCESSING":
                                print(f"[Backend] Waiting for file '{gemini_file.name}' to be processed...")
                                while state == "PROCESSING":
                                    time.sleep(2)
                                    gemini_file = client.files.get(name=gemini_file.name)
                                    state = getattr(gemini_file.state, "name", gemini_file.state) if hasattr(gemini_file, "state") else None
                                if state == "FAILED":
                                    raise Exception("Gemini failed to process the uploaded file.")
                                    
                            gemini_files.append(gemini_file)
                            print(f"[Backend] Successfully uploaded '{uploaded.name}' as {gemini_file.name}")
                        finally:
                            try:
                                os.remove(temp_path)
                            except OSError:
                                # If deletion fails, just ignore; it will stay in OS temp dir.
                                pass

                combined_text = "\n\n".join(text_parts)

                try:
                    summary = summarize_notes(combined_text, gemini_files, selected_model, has_pdf)
                        
                    if not summary:
                        st.error("Gemini returned an empty summary. Please try again.")
                    else:
                        st.session_state.summary = summary
                        st.success("✅ Summary generated successfully!")
                        st.rerun()
                except Exception as e:
                    st.error(f"❌ Error while generating summary: {e}")
                finally:
                    # Remove the overlay loader in case of error
                    loader_placeholder.empty()
                    # Clean up uploaded files from Gemini servers to prevent storage quota exhaustion
                    for g_file in gemini_files:
                        try:
                            print(f"[Backend] Cleaning up: Deleting {g_file.name} from Gemini servers...")
                            client.files.delete(name=g_file.name)
                            print(f"[Backend] Successfully deleted {g_file.name}.")
                        except Exception as delete_err:
                            print(f"[Backend] WARNING: Failed to delete {g_file.name}: {delete_err}")

    # Show the summary if available
    if st.session_state.summary:
        st.subheader("📚 Summarized Notes")
        st.markdown(st.session_state.summary)

        # Action button
        with st.container():
            # Print / Save as PDF using browser (exact visual copy)
            if st.button("🖨️ Print / Save as PDF", key="print_summary"):
                # Open a minimal HTML document containing only the summarized notes,
                # then trigger the browser print dialog. This keeps formatting very
                # close to what is shown in the app while giving the browser full
                # control over PDF generation.
                summary_html = markdown(st.session_state.summary, extensions=["tables"])
                printable_html = f"""
                <html>
                  <head>
                    <meta charset="utf-8" />
                    <title>Summarized Notes</title>
                    <style>
                      /* Screen styles */
                      body {{
                        font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
                        padding: 24px;
                        background-color: #020617;
                        color: #e5e7eb;
                      }}
                      h1, h2, h3, h4 {{
                        color: #60a5fa;
                        margin-top: 1.4rem;
                        margin-bottom: 0.6rem;
                      }}
                      p {{
                        margin: 0.4rem 0;
                        line-height: 1.5;
                      }}
                      ul, ol {{
                        margin: 0.4rem 0 0.8rem 1.5rem;
                      }}
                      table {{
                        border-collapse: collapse;
                        width: 100%;
                        margin: 1rem 0;
                      }}
                      th, td {{
                        border: 1px solid #4b5563;
                        padding: 0.4rem 0.6rem;
                      }}
                      th {{
                        background-color: #111827;
                      }}
                      
                      /* Print/PDF styles - Remove headers, footers, and overlays */
                      @media print {{
                        @page {{
                          margin: 0.75in;
                          size: auto;
                        }}
                        
                        /* Remove all browser default headers and footers */
                        @page {{
                          margin-top: 0.5in;
                          margin-bottom: 0.5in;
                        }}
                        
                        /* Hide any overlays or loaders */
                        .overlay-loader,
                        .overlay-loader-circle,
                        .overlay-loader-text,
                        header,
                        footer {{
                          display: none !important;
                          visibility: hidden !important;
                        }}
                        
                        /* Remove blur effects */
                        * {{
                          backdrop-filter: none !important;
                          -webkit-backdrop-filter: none !important;
                          filter: none !important;
                          -webkit-filter: none !important;
                        }}
                        
                        /* Clear, readable content */
                        body {{
                          background: white !important;
                          color: #000000 !important;
                          padding: 0 !important;
                          -webkit-print-color-adjust: exact;
                          print-color-adjust: exact;
                        }}
                        
                        /* Black text for better readability */
                        h1, h2, h3, h4, h5, h6 {{
                          color: #000000 !important;
                          page-break-after: avoid;
                        }}
                        
                        p, li, td, th, span, div {{
                          color: #000000 !important;
                        }}
                        
                        /* Remove shadows that might cause blur */
                        * {{
                          box-shadow: none !important;
                          text-shadow: none !important;
                        }}
                        
                        /* Ensure tables print well */
                        table {{
                          border-collapse: collapse;
                          page-break-inside: avoid;
                        }}
                        
                        th, td {{
                          border: 1px solid #000000 !important;
                          padding: 0.4rem 0.6rem;
                        }}
                        
                        th {{
                          background-color: #f0f0f0 !important;
                          color: #000000 !important;
                        }}
                      }}
                    </style>
                  </head>
                  <body>
                    {summary_html}
                    <script>
                      window.onload = function() {{ window.print(); }};
                    </script>
                  </body>
                </html>
                """
                components.html(printable_html, height=0, width=0)

    # Render footer
    render_footer()


if __name__ == "__main__":
    main()